#!/usr/bin/env python3
"""
Apex Sports Club — Enhanced Presentation Generator v2
Embeds live UI screenshots + database schema slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────
PRIMARY      = RGBColor(0x25, 0x63, 0xEB)
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)
DARK_CARD    = RGBColor(0x1E, 0x29, 0x3B)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_PURPLE= RGBColor(0x7C, 0x3A, 0xED)
ACCENT_ORANGE= RGBColor(0xF9, 0x73, 0x16)
ACCENT_RED   = RGBColor(0xEF, 0x44, 0x44)
ACCENT_CYAN  = RGBColor(0x06, 0xB6, 0xD4)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS_DIR = os.path.join(SCRIPT_DIR, '..', 'screenshots')


def set_slide_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(tf, text, font_size=14, color=WHITE, bold=False, space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

def add_badge(slide, left, top, text, bg_color, text_color=WHITE, font_size=11):
    w = Inches(len(text) * 0.1 + 0.4)
    h = Inches(0.35)
    shape = add_shape_rect(slide, left, top, w, h, bg_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(font_size)
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Calibri"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.word_wrap = False
    return shape

def ui_card(slide, left, top, width, height, title, items, accent_color=PRIMARY):
    card = add_shape_rect(slide, left, top, width, height, DARK_CARD, RGBColor(0x33, 0x40, 0x55))
    bar = add_shape_rect(slide, left, top, width, Inches(0.45), accent_color)
    bar.text_frame.paragraphs[0].text = title
    bar.text_frame.paragraphs[0].font.size = Pt(12)
    bar.text_frame.paragraphs[0].font.color.rgb = WHITE
    bar.text_frame.paragraphs[0].font.bold = True
    bar.text_frame.paragraphs[0].font.name = "Calibri"
    bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    txBox = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.6), width - Inches(0.5), height - Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            tf.paragraphs[0].text = f"*  {item}"
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
            tf.paragraphs[0].font.name = "Calibri"
            tf.paragraphs[0].space_after = Pt(4)
        else:
            add_paragraph(tf, f"*  {item}", font_size=11, color=LIGHT_GRAY, space_before=2, space_after=2)
    return card

def add_screenshot(slide, img_name, left, top, width, height, label=""):
    """Embed a screenshot image into a slide."""
    path = os.path.join(SCREENSHOTS_DIR, f"{img_name}.png")
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
        if label:
            add_text_box(slide, left, top + height + Inches(0.05), width, Inches(0.3),
                         label, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        return True
    return False

def db_table_card(slide, left, top, width, height, table_name, columns, accent_color=PRIMARY):
    """Draw a compact database table card."""
    card = add_shape_rect(slide, left, top, width, height, DARK_CARD, accent_color)
    # Table name header
    header = add_shape_rect(slide, left, top, width, Inches(0.35), accent_color)
    header.text_frame.paragraphs[0].text = table_name
    header.text_frame.paragraphs[0].font.size = Pt(10)
    header.text_frame.paragraphs[0].font.color.rgb = WHITE
    header.text_frame.paragraphs[0].font.bold = True
    header.text_frame.paragraphs[0].font.name = "Consolas"
    header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Columns
    txBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.4), width - Inches(0.2), height - Inches(0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, col in enumerate(columns):
        if i == 0:
            tf.paragraphs[0].text = col
            tf.paragraphs[0].font.size = Pt(8)
            tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
            tf.paragraphs[0].font.name = "Consolas"
            tf.paragraphs[0].space_after = Pt(1)
        else:
            add_paragraph(tf, col, font_size=8, color=LIGHT_GRAY, space_before=1, space_after=1)

# ══════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

# ──────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY)
badge = add_shape_rect(slide, Inches(1), Inches(1.5), Inches(4), Inches(0.6), PRIMARY)
badge.text_frame.paragraphs[0].text = "APEX SPORTS CLUB"
badge.text_frame.paragraphs[0].font.size = Pt(14)
badge.text_frame.paragraphs[0].font.color.rgb = WHITE
badge.text_frame.paragraphs[0].font.bold = True
badge.text_frame.paragraphs[0].font.name = "Calibri"
badge.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
add_text_box(slide, Inches(1), Inches(2.4), Inches(10), Inches(1.5),
             "Apex Sports Club", font_size=52, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.8), Inches(10), Inches(1),
             "Comprehensive Sports Management Platform", font_size=24, color=LIGHT_GRAY)
add_text_box(slide, Inches(1), Inches(5.0), Inches(10), Inches(0.6),
             "PHP | MySQL | Bootstrap 5 | Cloudflare Turnstile | Paystack | Brevo Email | Gemini AI",
             font_size=14, color=RGBColor(0x64, 0x74, 0x8B))
add_text_box(slide, Inches(1), Inches(6.2), Inches(6), Inches(0.4),
             "System Presentation  |  July 2026", font_size=13, color=LIGHT_GRAY)

# ──────────────────────────────────────────────────────────────
# SLIDE 2 — Live UI: Homepage & Login
# ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Live UI — Homepage & Login", font_size=36, color=WHITE, bold=True)
add_badge(slide, Inches(9), Inches(0.45), "LIVE SCREENSHOTS", ACCENT_GREEN)

add_screenshot(slide, "homepage", Inches(0.8), Inches(1.4), Inches(6), Inches(3.375), "Homepage — Hero section with 3D background")
add_screenshot(slide, "login", Inches(7.1), Inches(1.4), Inches(6), Inches(3.375), "Login Page — Secure authentication")

ui_card(slide, Inches(0.8), Inches(5.2), Inches(5.8), Inches(1.8),
        "Key Features", [
            "Modern dark-themed UI with 3D Spline robot background",
            "Club stats: 1,200+ members, 24/7 access, 12+ leagues",
            "Cloudflare Turnstile CAPTCHA integration",
            "Responsive design with mobile-first approach",
        ], PRIMARY)
ui_card(slide, Inches(7.0), Inches(5.2), Inches(5.8), Inches(1.8),
        "Auth Features", [
            "Session-based authentication with CSRF protection",
            "Password strength policy enforcement",
            "Rate limiting: max 3 attempts per IP/hour",
            "TOTP 2FA support for admin accounts",
        ], ACCENT_RED)

# ──────────────────────────────────────────────────────────────
# SLIDE 3 — Live UI: Registration
# ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Live UI — Registration", font_size=36, color=WHITE, bold=True)
add_badge(slide, Inches(7.5), Inches(0.45), "ONBOARDING", ACCENT_GREEN)

add_screenshot(slide, "register", Inches(1.5), Inches(1.3), Inches(10.3), Inches(5.8))

# ──────────────────────────────────────────────────────────────
# SLIDE 4 — Live UI: Public Pages Grid
# ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Live UI — Public Pages", font_size=36, color=WHITE, bold=True)
add_badge(slide, Inches(8), Inches(0.45), "4 PAGES", ACCENT_ORANGE)

# 2x2 grid of screenshots
screenshots_grid = [
    ("view_sports",     Inches(0.8),  Inches(1.3), "Browse Sports"),
    ("view_facilities", Inches(6.9),  Inches(1.3), "View Facilities"),
    ("view_fixtures",   Inches(0.8),  Inches(4.2), "Fixtures & Standings"),
    ("booking",         Inches(6.9),  Inches(4.2), "Book a Session"),
]
for img_name, x, y, label in screenshots_grid:
    add_screenshot(slide, img_name, x, y, Inches(5.7), Inches(2.55), label)

# ──────────────────────────────────────────────────────────────
# SLIDE 5 — Live UI: Admin Dashboard
# ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_RED)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Live UI — Admin Panel", font_size=36, color=WHITE, bold=True)
add_badge(slide, Inches(7), Inches(0.45), "ADMIN", ACCENT_RED)

add_screenshot(slide, "admin_login", Inches(0.8), Inches(1.3), Inches(5.8), Inches(2.8), "Admin Login")
add_screenshot(slide, "admin_dashboard", Inches(7.0), Inches(1.3), Inches(5.8), Inches(2.8), "Admin Dashboard")

ui_card(slide, Inches(0.8), Inches(4.6), Inches(12), Inches(2.5),
        "Admin Capabilities", [
            "Real-time stats: Members, Bookings, Revenue, Leagues — all live from the database",
            "AI Booking Review: automated approval/rejection with configurable strictness levels",
            "35+ management pages: Members, Sports, Coaches, Facilities, Equipment, Payments, Leagues, Teams, Fixtures",
            "Smart features: Churn prediction, AI match reports (Gemini), Smart scheduling, Bulk email campaigns",
        ], ACCENT_RED)

# ──────────────────────────────────────────────────────────────
# SLIDE 6 — Database Schema
# ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Database Schema", font_size=36, color=WHITE, bold=True)
add_badge(slide, Inches(7.5), Inches(0.45), "49 MIGRATIONS", ACCENT_CYAN)

# Core tables
db_table_card(slide, Inches(0.5), Inches(1.3), Inches(2.4), Inches(2.6),
    "members", [
        "member_id      INT PK",
        "first_name     VARCHAR",
        "last_name      VARCHAR",
        "email          VARCHAR UNQ",
        "password       VARCHAR",
        "phone          VARCHAR",
        "referral_code  VARCHAR",
        "role           ENUM",
    ], PRIMARY)

db_table_card(slide, Inches(3.2), Inches(1.3), Inches(2.4), Inches(2.6),
    "bookings", [
        "booking_id     INT PK",
        "member_id      INT FK",
        "facility_id    INT FK",
        "coach_id       INT FK",
        "sport_id       INT FK",
        "booking_date   DATE",
        "start_time     TIME",
        "status         ENUM",
    ], ACCENT_ORANGE)

db_table_card(slide, Inches(5.9), Inches(1.3), Inches(2.4), Inches(2.6),
    "payments", [
        "payment_id     INT PK",
        "member_id      INT FK",
        "booking_id     INT FK",
        "amount         DECIMAL",
        "payment_method VARCHAR",
        "payment_date   DATETIME",
        "transaction_id VARCHAR",
        "status         ENUM",
    ], ACCENT_GREEN)

db_table_card(slide, Inches(8.6), Inches(1.3), Inches(2.4), Inches(2.6),
    "leagues", [
        "league_id      INT PK",
        "name           VARCHAR",
        "sport_id       INT FK",
        "season         VARCHAR",
        "status         ENUM",
        "start_date     DATE",
        "end_date       DATE",
    ], ACCENT_PURPLE)

db_table_card(slide, Inches(11.3), Inches(1.3), Inches(2.4), Inches(2.6),
    "fixtures", [
        "fixture_id     INT PK",
        "league_id      INT FK",
        "home_team_id   INT FK",
        "away_team_id   INT FK",
        "match_date     DATE",
        "home_score     INT",
        "away_score     INT",
        "status         ENUM",
    ], ACCENT_RED)

# Second row
db_table_card(slide, Inches(0.5), Inches(4.2), Inches(2.4), Inches(2.6),
    "sports", [
        "sport_id       INT PK",
        "name           VARCHAR",
        "description    TEXT",
        "icon           VARCHAR",
    ], PRIMARY)

db_table_card(slide, Inches(3.2), Inches(4.2), Inches(2.4), Inches(2.6),
    "coaches", [
        "coach_id       INT PK",
        "first_name     VARCHAR",
        "last_name      VARCHAR",
        "specialization VARCHAR",
        "hourly_rate    DECIMAL",
    ], ACCENT_ORANGE)

db_table_card(slide, Inches(5.9), Inches(4.2), Inches(2.4), Inches(2.6),
    "facilities", [
        "facility_id    INT PK",
        "name           VARCHAR",
        "sport_id       INT FK",
        "capacity       INT",
        "hourly_rate    DECIMAL",
        "status         ENUM",
    ], ACCENT_GREEN)

db_table_card(slide, Inches(8.6), Inches(4.2), Inches(2.4), Inches(2.6),
    "teams", [
        "team_id        INT PK",
        "name           VARCHAR",
        "league_id      INT FK",
        "logo_url       VARCHAR",
        "manager_name   VARCHAR",
    ], ACCENT_PURPLE)

db_table_card(slide, Inches(11.3), Inches(4.2), Inches(2.4), Inches(2.6),
    "admins", [
        "admin_id       INT PK",
        "email          VARCHAR UNQ",
        "password       VARCHAR",
        "role           ENUM",
        "totp_secret    VARCHAR",
        "is_2fa_enabled TINYINT",
    ], ACCENT_RED)

# Schema summary
add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
             "Total: 60+ tables  |  49 migration files  |  Foreign key relationships across all core entities",
             font_size=13, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────────
# SLIDE 7 — Feature Overview
# ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Platform Features", font_size=36, color=WHITE, bold=True)

features_grid = [
    ("MEMBERSHIP", PRIMARY, [
        "Member registration & login",
        "Digital membership cards (QR)",
        "Referral & loyalty system",
        "Membership pause & renewal",
        "Parent portal & compliance",
        "GDPR data export",
    ]),
    ("BOOKING", ACCENT_ORANGE, [
        "Facility & coach booking",
        "Booking calendar view",
        "AI automated review",
        "Smart scheduling suggestions",
        "Attendance tracking",
        "Session notes & ratings",
    ]),
    ("COMPETITION", ACCENT_GREEN, [
        "League & team management",
        "Fixture scheduling",
        "Live scoring & match events",
        "Standings auto-recalculation",
        "Top scorers leaderboard",
        "MOTM voting",
    ]),
    ("FINANCE", ACCENT_PURPLE, [
        "Paystack checkout",
        "M-Pesa Daraja integration",
        "Revenue dashboard",
        "Refund management",
        "Promo codes",
        "Membership billing",
    ]),
    ("AI / ML", ACCENT_CYAN, [
        "Booking review automation",
        "Churn prediction",
        "AI match reports (Gemini)",
        "Insights engine",
        "Custom AI prompts",
        "Multi-model fallback",
    ]),
    ("ENGAGEMENT", ACCENT_RED, [
        "Forum & fan wall",
        "Polls & announcements",
        "Gallery & highlight reels",
        "Gear marketplace",
        "Volunteer management",
        "WhatsApp widget",
    ]),
]

for i, (title, color, items) in enumerate(features_grid):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.3 + row * 3.1)
    ui_card(slide, x, y, Inches(3.8), Inches(2.9), title, items, color)

# ──────────────────────────────────────────────────────────────
# SLIDE 8 — Tech Stack & Architecture
# ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Technical Architecture", font_size=36, color=WHITE, bold=True)

tech_items = [
    ("Backend", "PHP 7.4+\nMySQL / MariaDB\nmysqli extension\nSession management", PRIMARY),
    ("Frontend", "Bootstrap 5\nVanilla JavaScript\nGoogle Fonts (Inter)\nFont Awesome 6", ACCENT_GREEN),
    ("Integrations", "Paystack payments\nM-Pesa Daraja\nBrevo email\nCloudflare Turnstile", ACCENT_ORANGE),
    ("AI / ML", "Google Gemini API\nOpenRouter (200+ models)\nChurn prediction\nSmart scheduling", ACCENT_PURPLE),
]

for i, (title, items, color) in enumerate(tech_items):
    x = Inches(0.8 + i * 3.1)
    card = add_shape_rect(slide, x, Inches(1.4), Inches(2.8), Inches(2.4), DARK_CARD, color)
    add_text_box(slide, x + Inches(0.2), Inches(1.5), Inches(2.4), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(2.4), Inches(1.6),
                 items, font_size=13, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(4.1), Inches(11), Inches(0.5),
             "Project Structure", font_size=20, color=WHITE, bold=True)

structure = [
    ("admin/", "35+ admin pages — dashboard, CRUD, AI tools"),
    ("public/", "40+ member pages — booking, profile, leagues"),
    ("includes/", "30+ shared modules — auth, email, payments, AI"),
    ("config/", "Database & API configuration loaders"),
    ("callbacks/", "Payment callback endpoints (Paystack, M-Pesa)"),
    ("migrations/", "49 SQL migration files — full schema"),
    ("scripts/", "Utility scripts — roster generation, testing"),
    ("dev/", "Development & testing utilities"),
]

for i, (folder, desc) in enumerate(structure):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(4.7 + row * 0.65)
    add_text_box(slide, x, y, Inches(1.5), Inches(0.4),
                 folder, font_size=13, color=PRIMARY, bold=True)
    add_text_box(slide, x + Inches(1.5), y, Inches(4.5), Inches(0.4),
                 desc, font_size=12, color=LIGHT_GRAY)

# ──────────────────────────────────────────────────────────────
# SLIDE 9 — Thank You
# ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY)

badge = add_shape_rect(slide, Inches(5.4), Inches(1.2), Inches(2.5), Inches(0.6), PRIMARY)
badge.text_frame.paragraphs[0].text = "APEX SPORTS CLUB"
badge.text_frame.paragraphs[0].font.size = Pt(12)
badge.text_frame.paragraphs[0].font.color.rgb = WHITE
badge.text_frame.paragraphs[0].font.bold = True
badge.text_frame.paragraphs[0].font.name = "Calibri"
badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(2), Inches(2.2), Inches(9.3), Inches(1.2),
             "Thank You", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(3.5), Inches(9.3), Inches(0.8),
             "Apex Sports Club — Built with passion for the game",
             font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

features = [
    "Real-time Booking", "AI-Powered Insights", "Secure Payments",
    "League Management", "Mobile Responsive", "Enterprise Security"
]
for i, feat in enumerate(features):
    x = Inches(1.5 + (i % 3) * 3.6)
    y = Inches(4.8 + (i // 3) * 0.7)
    add_text_box(slide, x, y, Inches(3.2), Inches(0.5),
                 feat, font_size=15, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(6.4), Inches(9.3), Inches(0.5),
             "PHP | MySQL | Bootstrap 5 | Cloudflare | Paystack | Brevo | Gemini AI",
             font_size=13, color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
if __name__ == '__main__':
    output_path = os.path.join(SCRIPT_DIR, '..', 'Apex_Sports_Club_System_Presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved to: {os.path.abspath(output_path)}")
    print(f"Slides: {len(prs.slides)}")
